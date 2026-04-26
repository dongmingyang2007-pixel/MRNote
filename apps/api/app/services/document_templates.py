"""Generate blank Office / PDF documents.

Pure byte producers — no DB, no S3, no external services. Callers persist
the bytes and create the matching DataItem/StudyAsset rows themselves.
"""

from __future__ import annotations

import io
from typing import Literal

from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer


DocType = Literal["docx", "xlsx", "pptx", "pdf"]


# Maps doc_type → (media_type, default extension, study asset_type)
DOC_TYPE_MEDIA: dict[DocType, tuple[str, str, str]] = {
    "docx": (
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "docx",
        "pdf",  # study pipeline treats word docs as document-style assets
    ),
    "xlsx": (
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "xlsx",
        "article",
    ),
    "pptx": (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "pptx",
        "slides",
    ),
    "pdf": ("application/pdf", "pdf", "pdf"),
}


def create_blank_docx(title: str = "Untitled") -> bytes:
    """One-paragraph .docx ready for ONLYOFFICE to open."""
    document = Document()
    document.core_properties.title = title or "Untitled"
    document.add_paragraph(title or "")
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def create_blank_xlsx(title: str = "Untitled") -> bytes:
    """One-sheet .xlsx with the title on Sheet1!A1 only when non-empty."""
    workbook = Workbook()
    worksheet = workbook.active
    worksheet.title = "Sheet1"
    if title:
        worksheet["A1"] = title
    workbook.properties.title = title or "Untitled"
    buffer = io.BytesIO()
    workbook.save(buffer)
    return buffer.getvalue()


def create_blank_pptx(title: str = "Untitled") -> bytes:
    """Single-slide .pptx with title placeholder filled in."""
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    title_layout = presentation.slide_layouts[0]
    slide = presentation.slides.add_slide(title_layout)
    if slide.shapes.title is not None:
        slide.shapes.title.text = title or "Untitled"
        for paragraph in slide.shapes.title.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(40)
    presentation.core_properties.title = title or "Untitled"
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def create_blank_pdf(title: str = "Untitled") -> bytes:
    """One-page A4 PDF with the title centered near the top."""
    buffer = io.BytesIO()
    document = SimpleDocTemplate(
        buffer,
        pagesize=A4,
        title=title or "Untitled",
        leftMargin=72,
        rightMargin=72,
        topMargin=72,
        bottomMargin=72,
    )
    styles = getSampleStyleSheet()
    document.build(
        [
            Spacer(1, 36),
            Paragraph(title or "Untitled", styles["Title"]),
        ]
    )
    return buffer.getvalue()


_GENERATORS = {
    "docx": create_blank_docx,
    "xlsx": create_blank_xlsx,
    "pptx": create_blank_pptx,
    "pdf": create_blank_pdf,
}


def generate_blank_document(doc_type: DocType, title: str = "Untitled") -> bytes:
    if doc_type not in _GENERATORS:
        raise ValueError(f"Unsupported doc_type: {doc_type}")
    return _GENERATORS[doc_type](title)
